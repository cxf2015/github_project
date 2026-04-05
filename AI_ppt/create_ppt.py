#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
为二年级学生制作"AI不是看听说"科普PPT
卡通风格、简单易懂、10页
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿（16:9宽屏）
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色方案（卡通明亮风格）
COLORS = {
    'primary_blue': RGBColor(79, 172, 254),      # 天蓝色
    'primary_green': RGBColor(67, 233, 123),     # 草绿色
    'accent_yellow': RGBColor(255, 214, 102),    # 明黄色
    'accent_pink': RGBColor(255, 154, 158),      # 柔粉色
    'accent_purple': RGBColor(168, 139, 235),    # 浅紫色
    'white': RGBColor(255, 255, 255),
    'dark_text': RGBColor(51, 51, 51),
    'light_bg': RGBColor(245, 250, 255),
}

# 辅助函数：添加圆角矩形背景
def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # 将形状移到最底层
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape

# 辅助函数：添加装饰圆形
def add_circle(slide, left, top, size, fill_color, transparency=0.3):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# 辅助函数：添加文本框
def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=None, align=PP_ALIGN.CENTER, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color if color else COLORS['dark_text']
    return txBox

# 辅助函数：添加标题
def add_title(slide, text, top=Inches(0.4)):
    return add_textbox(slide, Inches(0.5), top, Inches(12.333), Inches(1),
                       text, 44, bold=True, color=COLORS['primary_blue'], align=PP_ALIGN.CENTER)

# 辅助函数：添加内容文本
def add_content(slide, text, top, left=Inches(1), width=Inches(11.333), font_size=24):
    return add_textbox(slide, left, top, width, Inches(1.5),
                       text, font_size, color=COLORS['dark_text'], align=PP_ALIGN.LEFT)

# ==================== 第1页：封面 ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

# 背景装饰圆
add_circle(slide1, Inches(-1), Inches(-1), Inches(4), COLORS['accent_yellow'], 0.2)
add_circle(slide1, Inches(10), Inches(-0.5), Inches(3), COLORS['accent_pink'], 0.2)
add_circle(slide1, Inches(-0.5), Inches(5), Inches(2.5), COLORS['primary_green'], 0.2)
add_circle(slide1, Inches(11), Inches(5.5), Inches(2), COLORS['accent_purple'], 0.2)

# 主标题背景
add_rounded_rect(slide1, Inches(1.5), Inches(2), Inches(10.333), Inches(2), COLORS['white'])

# 主标题
add_textbox(slide1, Inches(1), Inches(2.2), Inches(11.333), Inches(1.2),
            "AI不是看听说", 60, bold=True, color=COLORS['primary_blue'], align=PP_ALIGN.CENTER)

# 副标题
add_textbox(slide1, Inches(1), Inches(3.3), Inches(11.333), Inches(0.8),
            "—— 揭秘AI的小秘密 ——", 32, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

# AI机器人图标（用形状组合表示）
# 头部
head = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(4.3), Inches(1.8), Inches(1.5))
head.fill.solid()
head.fill.fore_color.rgb = COLORS['primary_blue']
head.line.fill.background()

# 眼睛
for i, x in enumerate([6.1, 6.5]):
    eye = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(4.6), Inches(0.25), Inches(0.25))
    eye.fill.solid()
    eye.fill.fore_color.rgb = COLORS['white']
    eye.line.fill.background()

# 天线
antenna = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(4), Inches(0.15), Inches(0.4))
antenna.fill.solid()
antenna.fill.fore_color.rgb = COLORS['accent_yellow']
top_ball = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.45), Inches(3.7), Inches(0.35), Inches(0.35))
top_ball.fill.solid()
top_ball.fill.fore_color.rgb = COLORS['accent_yellow']

# 身体
body = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(5.8), Inches(1.6), Inches(1))
body.fill.solid()
body.fill.fore_color.rgb = COLORS['primary_blue']
body.line.fill.background()

# 心形
try:
    heart = slide1.shapes.add_shape(MSO_SHAPE.HEART, Inches(6.4), Inches(6), Inches(0.6), Inches(0.6))
    heart.fill.solid()
    heart.fill.fore_color.rgb = COLORS['accent_pink']
    heart.line.fill.background()
except:
    pass  # 如果没有心形就用圆形

# 底部文字
add_textbox(slide1, Inches(1), Inches(6.8), Inches(11.333), Inches(0.5),
            "给二年级小朋友的人工智能科普课", 20, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

# ==================== 第2页：什么是AI ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# 装饰
add_circle(slide2, Inches(-0.5), Inches(-0.5), Inches(2), COLORS['primary_blue'], 0.1)
add_circle(slide2, Inches(11.5), Inches(6), Inches(2.5), COLORS['accent_yellow'], 0.1)

add_title(slide2, "什么是AI？")

# 内容框
add_rounded_rect(slide2, Inches(1.5), Inches(1.5), Inches(10.333), Inches(3.5), COLORS['light_bg'])

# 主要解释
add_content(slide2, "AI就是【聪明的电脑程序】", Inches(2), font_size=32)
add_content(slide2, "它会学习很多知识，帮我们解决问题", Inches(2.8), font_size=28)
add_content(slide2, "但它不是真正的人，只是工具哦！", Inches(3.6), font_size=28)

# 电脑卡通图标
# 显示器
monitor = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(5.2), Inches(2.8), Inches(1.8))
monitor.fill.solid()
monitor.fill.fore_color.rgb = COLORS['primary_blue']
monitor.line.fill.background()

# 屏幕
screen = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(5.4), Inches(2.4), Inches(1.4))
screen.fill.solid()
screen.fill.fore_color.rgb = COLORS['white']

# 屏幕上的AI文字
add_textbox(slide2, Inches(5.5), Inches(5.8), Inches(2.4), Inches(0.6),
            "AI", 28, bold=True, color=COLORS['primary_blue'], align=PP_ALIGN.CENTER)

# 底座
base = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.3), Inches(7), Inches(0.8), Inches(0.2))
base.fill.solid()
base.fill.fore_color.rgb = COLORS['dark_text']

# 支架
stand = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(6.95), Inches(0.3), Inches(0.15))
stand.fill.solid()
stand.fill.fore_color.rgb = COLORS['dark_text']

# ==================== 第3页：AI的"看" ====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

add_circle(slide3, Inches(11), Inches(-0.5), Inches(2.5), COLORS['accent_pink'], 0.15)
add_circle(slide3, Inches(-0.5), Inches(5.5), Inches(2), COLORS['primary_green'], 0.1)

add_title(slide3, 'AI的"看"——其实是"认像素"')

# 左侧说明框
add_rounded_rect(slide3, Inches(0.8), Inches(1.4), Inches(5.5), Inches(2.8), COLORS['light_bg'])
add_content(slide3, '• AI把图片分成很多小格子', Inches(1.6), left=Inches(1), width=Inches(5))
add_content(slide3, '• 每个格子叫"像素"', Inches(2.3), left=Inches(1), width=Inches(5))
add_content(slide3, '• AI通过颜色判断是什么', Inches(3), left=Inches(1), width=Inches(5))

# 右侧例子框
add_rounded_rect(slide3, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.8), COLORS['accent_yellow'])
add_textbox(slide3, Inches(7), Inches(1.6), Inches(5.4), Inches(0.6),
            "例子", 26, bold=True, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)
add_content(slide3, "• 手机拍照识别植物", Inches(2.4), left=Inches(7.2), width=Inches(5))
add_content(slide3, "• 找相似的图片", Inches(3), left=Inches(7.2), width=Inches(5))
add_content(slide3, "• 人脸识别解锁手机", Inches(3.6), left=Inches(7.2), width=Inches(5))

# 像素图示（简化的网格）
grid_y = Inches(5)
for row in range(4):
    for col in range(4):
        x = Inches(3 + col * 0.4)
        y = grid_y + Inches(row * 0.35)
        pixel = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.35), Inches(0.3))
        # 不同颜色代表不同像素
        colors = [COLORS['primary_blue'], COLORS['accent_yellow'], COLORS['accent_pink'], COLORS['primary_green']]
        pixel.fill.solid()
        pixel.fill.fore_color.rgb = colors[(row + col) % 4]
        pixel.line.color.rgb = COLORS['white']

add_textbox(slide3, Inches(2), Inches(6.5), Inches(5), Inches(0.5),
            "就像这样，一格一格地看", 20, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

# 底部提示
add_textbox(slide3, Inches(6.5), Inches(5.5), Inches(6.5), Inches(1.5),
            'AI不会真正"看见"，\n它只是在做数学计算！', 24,
            color=COLORS['primary_blue'], align=PP_ALIGN.CENTER)

# ==================== 第4页：AI的"听" ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

add_circle(slide4, Inches(-0.3), Inches(-0.3), Inches(1.8), COLORS['accent_purple'], 0.1)
add_circle(slide4, Inches(11.5), Inches(5), Inches(3), COLORS['primary_blue'], 0.1)

add_title(slide4, 'AI的"听"——其实是"转文字"')

# 中间大框
add_rounded_rect(slide4, Inches(1.5), Inches(1.4), Inches(10.333), Inches(2.5), COLORS['light_bg'])
add_content(slide4, "• AI把声音变成波形图", Inches(1.7), left=Inches(2), font_size=28)
add_content(slide4, "• 再把波形图转换成文字", Inches(2.4), left=Inches(2), font_size=28)
add_content(slide4, '• 就像在玩"猜声音"的游戏', Inches(3.1), left=Inches(2), font_size=28)

# 声波图示
wave_y = Inches(4.3)
for i in range(15):
    x = Inches(3 + i * 0.5)
    height = 0.3 + 0.4 * abs((i % 5) - 2.5) / 2.5
    bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, wave_y - Inches(height/2),
                                   Inches(0.25), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary_green'] if i % 2 == 0 else COLORS['primary_blue']
    bar.line.fill.background()

add_textbox(slide4, Inches(2.5), Inches(5.2), Inches(8), Inches(0.5),
            "声音波形", 20, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

# 例子框
add_rounded_rect(slide4, Inches(1), Inches(5.8), Inches(11.333), Inches(1.4), COLORS['accent_yellow'])
add_textbox(slide4, Inches(1.2), Inches(6), Inches(11), Inches(1),
            "例子：语音助手、听歌识曲、语音输入", 26,
            color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

# ==================== 第5页：AI的"说" ====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

add_circle(slide5, Inches(10.5), Inches(-0.5), Inches(3), COLORS['accent_yellow'], 0.15)
add_circle(slide5, Inches(0), Inches(5.5), Inches(2.5), COLORS['accent_pink'], 0.1)

add_title(slide5, 'AI的"说"——其实是"播录音"')

# 左侧说明
add_rounded_rect(slide5, Inches(0.8), Inches(1.4), Inches(5.5), Inches(3.2), COLORS['light_bg'])
add_content(slide5, '• AI不会真的"说话"', Inches(1.7), left=Inches(1.2), font_size=26)
add_content(slide5, "• 它是把文字变成声音", Inches(2.4), left=Inches(1.2), font_size=26)
add_content(slide5, "• 像播放录音一样读出来", Inches(3.1), left=Inches(1.2), font_size=26)
add_content(slide5, "• 没有真正的情感哦", Inches(3.8), left=Inches(1.2), font_size=26)

# 右侧例子
add_rounded_rect(slide5, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.2), COLORS['primary_green'])
add_textbox(slide5, Inches(7), Inches(1.6), Inches(5.4), Inches(0.6),
            "生活中的例子", 26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
add_textbox(slide5, Inches(7.2), Inches(2.4), Inches(5), Inches(2),
            "• 导航语音播报\n• 故事机讲故事\n• 学习机读课文\n• 机器人回答问题",
            24, color=COLORS['white'], align=PP_ALIGN.LEFT)

# 喇叭图标
speaker = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(5), Inches(2.5), Inches(2.3))
speaker.fill.solid()
speaker.fill.fore_color.rgb = COLORS['primary_blue']
speaker.line.fill.background()

# 声波
for i in range(3):
    arc = slide5.shapes.add_shape(MSO_SHAPE.ARC, Inches(8.3 + i*0.4), Inches(5.3 + i*0.15),
                                   Inches(1.5 - i*0.3), Inches(1.8 - i*0.3))
    arc.fill.background()
    arc.line.color.rgb = COLORS['accent_yellow']
    arc.line.width = Pt(4)

add_textbox(slide5, Inches(5), Inches(7.3), Inches(3.5), Inches(0.4),
            "声音", 20, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

# 底部提示
add_textbox(slide5, Inches(0.5), Inches(4.8), Inches(5), Inches(0.8),
            "AI的声音\n是合成的哦！", 22,
            color=COLORS['primary_blue'], align=PP_ALIGN.CENTER)

# ==================== 第6页：AI能帮我们做什么（上） ====================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

add_circle(slide6, Inches(-0.5), Inches(0), Inches(2), COLORS['accent_yellow'], 0.1)
add_circle(slide6, Inches(11), Inches(4), Inches(3), COLORS['primary_green'], 0.1)

add_title(slide6, "AI能帮我们做什么？(一)")

# 学习助手框
add_rounded_rect(slide6, Inches(0.8), Inches(1.4), Inches(5.8), Inches(2.8), RGBColor(227, 242, 253))
add_textbox(slide6, Inches(1), Inches(1.5), Inches(5.4), Inches(0.6),
            "学习助手", 30, bold=True, color=COLORS['primary_blue'], align=PP_ALIGN.CENTER)
add_textbox(slide6, Inches(1.2), Inches(2.3), Inches(5.2), Inches(1.8),
            "• 查资料、找答案\n• 学英语、练口语\n• 讲解不会的题目",
            24, color=COLORS['dark_text'], align=PP_ALIGN.LEFT)

# 生活帮手框
add_rounded_rect(slide6, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.8), RGBColor(232, 245, 233))
add_textbox(slide6, Inches(7), Inches(1.5), Inches(5.4), Inches(0.6),
            "生活帮手", 30, bold=True, color=COLORS['primary_green'], align=PP_ALIGN.CENTER)
add_textbox(slide6, Inches(7.2), Inches(2.3), Inches(5.2), Inches(1.8),
            "• 定闹钟、查天气\n• 找路、导航\n• 控制智能家居",
            24, color=COLORS['dark_text'], align=PP_ALIGN.LEFT)

# 书本图标
book = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.5), Inches(2), Inches(2.2))
book.fill.solid()
book.fill.fore_color.rgb = COLORS['accent_pink']
book.line.fill.background()

add_textbox(slide6, Inches(2.7), Inches(5.1), Inches(1.6), Inches(1),
            "书", 40, align=PP_ALIGN.CENTER)

# 房子图标
house = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(4.8), Inches(2.2), Inches(1.8))
house.fill.solid()
house.fill.fore_color.rgb = COLORS['accent_yellow']
house.line.fill.background()

# 屋顶
roof = slide6.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(8.3), Inches(4), Inches(2.6), Inches(1))
roof.fill.solid()
roof.fill.fore_color.rgb = COLORS['accent_yellow']
roof.line.fill.background()

# 门
door = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(5.6), Inches(0.8), Inches(1))
door.fill.solid()
door.fill.fore_color.rgb = COLORS['dark_text']

# ==================== 第7页：AI能帮我们做什么（下） ====================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

add_circle(slide7, Inches(10.5), Inches(-0.3), Inches(2.5), COLORS['accent_purple'], 0.15)
add_circle(slide7, Inches(0), Inches(4.5), Inches(2.5), COLORS['accent_yellow'], 0.1)

add_title(slide7, "AI能帮我们做什么？(二)")

# 创意伙伴框
add_rounded_rect(slide7, Inches(0.8), Inches(1.4), Inches(5.8), Inches(2.8), RGBColor(243, 229, 245))
add_textbox(slide7, Inches(1), Inches(1.5), Inches(5.4), Inches(0.6),
            "创意伙伴", 30, bold=True, color=COLORS['accent_purple'], align=PP_ALIGN.CENTER)
add_textbox(slide7, Inches(1.2), Inches(2.3), Inches(5.2), Inches(1.8),
            "• 帮忙画画\n• 一起编故事\n• 创作音乐和诗歌",
            24, color=COLORS['dark_text'], align=PP_ALIGN.LEFT)

# 游戏玩伴框
add_rounded_rect(slide7, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.8), RGBColor(255, 243, 224))
add_textbox(slide7, Inches(7), Inches(1.5), Inches(5.4), Inches(0.6),
            "游戏玩伴", 30, bold=True, color=RGBColor(255, 152, 0), align=PP_ALIGN.CENTER)
add_textbox(slide7, Inches(7.2), Inches(2.3), Inches(5.2), Inches(1.8),
            "• 下棋（围棋、象棋）\n• 猜谜语、玩成语接龙\n• 推荐好玩的游戏",
            24, color=COLORS['dark_text'], align=PP_ALIGN.LEFT)

# 调色板图标
palette = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.2), Inches(4.5), Inches(2.6), Inches(1.8))
palette.fill.solid()
palette.fill.fore_color.rgb = COLORS['accent_pink']
palette.line.fill.background()

# 颜料点
paint_colors = [COLORS['accent_yellow'], COLORS['primary_green'], COLORS['primary_blue']]
for i, color in enumerate(paint_colors):
    dot = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.6 + i*0.6), Inches(4.8), Inches(0.5), Inches(0.5))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

# 画笔
brush = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(5.2), Inches(1.2), Inches(0.2))
brush.fill.solid()
brush.fill.fore_color.rgb = COLORS['dark_text']
brush.rotation = -30

# 游戏手柄图标
controller = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(4.6), Inches(2.4), Inches(1.5))
controller.fill.solid()
controller.fill.fore_color.rgb = RGBColor(255, 152, 0)
controller.line.fill.background()

# 按键
for i, (x, y) in enumerate([(8.3, 4.9), (8.6, 5.1), (8.9, 4.9)]):
    btn = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.35), Inches(0.35))
    btn.fill.solid()
    btn.fill.fore_color.rgb = COLORS['white']
    btn.line.fill.background()

# ==================== 第8页：AI不是真正的人 ====================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])

add_circle(slide8, Inches(0.5), Inches(-0.5), Inches(2), COLORS['accent_pink'], 0.15)
add_circle(slide8, Inches(10), Inches(5.5), Inches(3), COLORS['primary_blue'], 0.1)

add_title(slide8, "重要提醒：AI不是真正的人！")

# 中央大框
add_rounded_rect(slide8, Inches(1.5), Inches(1.4), Inches(10.333), Inches(4), COLORS['light_bg'])

# 对比内容
add_textbox(slide8, Inches(1.8), Inches(1.6), Inches(4.5), Inches(3.5),
            "AI不会：\n\n• 真正思考\n• 有感情\n• 自己想做某事\n• 知道对错",
            26, color=COLORS['dark_text'], align=PP_ALIGN.LEFT)

add_textbox(slide8, Inches(6.5), Inches(1.6), Inches(4.5), Inches(3.5),
            "AI只是：\n\n• 按程序运行\n• 计算数据\n• 被人类设计\n• 执行指令的工具",
            26, color=COLORS['primary_blue'], align=PP_ALIGN.LEFT)

# 分隔线
line = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.3), Inches(1.6), Inches(0.05), Inches(3.5))
line.fill.solid()
line.fill.fore_color.rgb = COLORS['accent_yellow']
line.line.fill.background()

# 底部强调框
add_rounded_rect(slide8, Inches(2), Inches(5.6), Inches(9.333), Inches(1.4), COLORS['accent_yellow'])
add_textbox(slide8, Inches(2.2), Inches(5.9), Inches(9), Inches(0.9),
            "AI就像计算器、铅笔一样，是帮助我们的工具！",
            28, bold=True, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

# ==================== 第9页：正确使用AI ====================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])

add_circle(slide9, Inches(-0.3), Inches(0.3), Inches(1.8), COLORS['primary_green'], 0.1)
add_circle(slide9, Inches(11), Inches(4.5), Inches(2.5), COLORS['accent_pink'], 0.15)

add_title(slide9, "怎样正确使用AI？")

# 三个提示框
positions = [(1, 1.4, '要家长陪同'),
             (4.7, 1.4, '控制使用时间'),
             (8.4, 1.4, '自己也要思考')]

for x, y, title in positions:
    add_rounded_rect(slide9, Inches(x), Inches(y), Inches(3.5), Inches(2.5), COLORS['light_bg'])
    add_textbox(slide9, Inches(x), Inches(y+0.8), Inches(3.5), Inches(0.6),
                title, 26, bold=True, color=COLORS['primary_blue'], align=PP_ALIGN.CENTER)

# 详细说明
add_content(slide9, "• 有问题先问爸爸妈妈", Inches(4.2), left=Inches(2), font_size=24)
add_content(slide9, "• 不要完全依赖AI的答案", Inches(4.9), left=Inches(2), font_size=24)
add_content(slide9, "• 保护个人隐私，不随便告诉AI信息", Inches(5.6), left=Inches(2), font_size=24)

# 盾牌图标
shield = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(6.3), Inches(1.8), Inches(1))
shield.fill.solid()
shield.fill.fore_color.rgb = COLORS['primary_green']
shield.line.fill.background()

add_textbox(slide9, Inches(5.8), Inches(6.4), Inches(1.8), Inches(0.8),
            "安全", 24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

# 底部安全提示
add_textbox(slide9, Inches(0.5), Inches(6.2), Inches(5), Inches(1),
            "安全上网\n快乐学习！", 22, bold=True,
            color=COLORS['primary_green'], align=PP_ALIGN.CENTER)

# ==================== 第10页：结束页 ====================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])

# 装饰圆
add_circle(slide10, Inches(-1), Inches(-1), Inches(4), COLORS['accent_yellow'], 0.2)
add_circle(slide10, Inches(10), Inches(-0.5), Inches(3), COLORS['accent_pink'], 0.2)
add_circle(slide10, Inches(-0.5), Inches(5), Inches(2.5), COLORS['primary_green'], 0.2)
add_circle(slide10, Inches(11), Inches(5.5), Inches(2), COLORS['accent_purple'], 0.2)

# 主内容框
add_rounded_rect(slide10, Inches(1.5), Inches(1.8), Inches(10.333), Inches(3), COLORS['white'])

# 总结
add_textbox(slide10, Inches(1.5), Inches(2.2), Inches(10.333), Inches(1),
            "记住这句话", 32, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

add_textbox(slide10, Inches(1.5), Inches(3), Inches(10.333), Inches(1.2),
            '"AI是工具，不是人"', 48, bold=True,
            color=COLORS['primary_blue'], align=PP_ALIGN.CENTER)

# 感谢语
add_textbox(slide10, Inches(1.5), Inches(4.5), Inches(10.333), Inches(0.8),
            "谢谢小朋友们！", 36, color=COLORS['primary_green'], align=PP_ALIGN.CENTER)

# AI机器人挥手告别（简化的机器人）
# 头
head = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(5.3), Inches(1.3), Inches(1))
head.fill.solid()
head.fill.fore_color.rgb = COLORS['primary_blue']
head.line.fill.background()

# 眼睛
for x in [6.2, 6.7]:
    eye = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(5.55), Inches(0.2), Inches(0.2))
    eye.fill.solid()
    eye.fill.fore_color.rgb = COLORS['white']
    eye.line.fill.background()

# 微笑
smile = slide10.shapes.add_shape(MSO_SHAPE.ARC, Inches(6.2), Inches(5.8), Inches(0.9), Inches(0.4))
smile.fill.background()
smile.line.color.rgb = COLORS['white']
smile.line.width = Pt(3)

# 身体
body = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(6.3), Inches(1.1), Inches(0.8))
body.fill.solid()
body.fill.fore_color.rgb = COLORS['primary_blue']
body.line.fill.background()

# 挥手的手臂
arm = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(5.5), Inches(0.6), Inches(0.2))
arm.fill.solid()
arm.fill.fore_color.rgb = COLORS['primary_blue']
arm.rotation = -30

hand = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.6), Inches(5.2), Inches(0.3), Inches(0.3))
hand.fill.solid()
hand.fill.fore_color.rgb = COLORS['accent_yellow']
hand.line.fill.background()

# 再见文字
add_textbox(slide10, Inches(5.8), Inches(7.2), Inches(1.8), Inches(0.4),
            "拜拜！", 22, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

# 保存文件
output_path = "D:/Lenovo/github/github_project/AI_ppt/AI不是看听说_二年级科普.pptx"
prs.save(output_path)
print(f"PPT已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
